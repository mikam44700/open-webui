#!/usr/bin/env python3
"""Atelier documents LunarIA — outil de Max (SPEC-agent-documents).

Fabrique des documents d'entreprise FINIS à partir de contenu structuré, puis les
publie dans l'application via l'API officielle des Fichiers (pont Fichiers, même
modèle que le pont Notes). L'agent fournit les faits ; le CLI garantit la forme
(mise en page sobre, totaux CALCULÉS — jamais fournis, donc jamais faux).

Commandes :
  xlsx    --titre T --data-file spec.json --sortie out.xlsx     tableur (KPI, totaux)
  docx    --titre T --md contenu.md --sortie out.docx           document Word
  pdf     --titre T --md contenu.md --sortie out.pdf            rapport PDF
  pptx    --spec spec.json --sortie out.pptx                    présentation sobre
  publier --fichier out.xlsx                                    téléverse + lien de téléchargement

(Les présentations de haute qualité passent par le MCP presenton — voir SKILL.md.)
Sortie : texte simple, code 0 si succès. Auth du pont : LUNARIA_INTERNAL_API_KEY (env).
"""

from __future__ import annotations

import argparse
import json
import mimetypes
import os
import sys
import urllib.error
import urllib.request
import uuid

DEFAULT_APP_URL = "http://localhost:8080"
TIMEOUT_SECONDS = 60
BLEU_NUIT = "1F3864"  # accent sobre commun à tous les documents


def _fail(message: str) -> None:
    print(f"ERREUR : {message}", file=sys.stderr)
    sys.exit(1)


def _lire(path: str) -> str:
    try:
        with open(path, encoding="utf-8") as handle:
            return handle.read()
    except OSError as exc:
        _fail(f"Fichier illisible : {exc}")
    return ""


def _lire_json(path: str) -> dict:
    try:
        return json.loads(_lire(path))
    except json.JSONDecodeError as exc:
        _fail(f"JSON invalide ({path}) : {exc}")
    return {}


# ── xlsx : tableur sobre, totaux calculés ────────────────────────────────────
# spec.json : {"feuilles":[{"nom": "...", "colonnes": [{"nom": "...", "format": "euro"|"nombre"|"texte"}],
#              "lignes": [[...], ...], "total": true|false}]}


def cmd_xlsx(args: argparse.Namespace) -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font, PatternFill
    from openpyxl.utils import get_column_letter

    spec = _lire_json(args.data_file)
    feuilles = spec.get("feuilles") or []
    if not feuilles:
        _fail("spec sans feuilles : {\"feuilles\": [{\"nom\", \"colonnes\", \"lignes\"}]}")
    wb = Workbook()
    wb.remove(wb.active)
    entete_font = Font(bold=True, color="FFFFFF")
    entete_fill = PatternFill("solid", fgColor=BLEU_NUIT)
    for feuille in feuilles:
        ws = wb.create_sheet((feuille.get("nom") or "Feuille")[:31])
        colonnes = feuille.get("colonnes") or []
        lignes = feuille.get("lignes") or []
        ws.append([c.get("nom", "") for c in colonnes])
        for cellule in ws[1]:
            cellule.font, cellule.fill = entete_font, entete_fill
            cellule.alignment = Alignment(horizontal="center")
        for ligne in lignes:
            ws.append(list(ligne))
        formats = {"euro": "#,##0.00 €", "nombre": "#,##0.00"}
        for idx, col in enumerate(colonnes, start=1):
            fmt = formats.get(str(col.get("format", "texte")))
            largeur = max([len(str(col.get("nom", "")))] + [len(str(l[idx - 1])) for l in lignes if len(l) >= idx])
            ws.column_dimensions[get_column_letter(idx)].width = min(max(largeur + 4, 12), 48)
            if fmt:
                for rang in range(2, len(lignes) + 2):
                    ws.cell(row=rang, column=idx).number_format = fmt
        if feuille.get("total") and lignes:
            # Totaux CALCULÉS ici (jamais fournis par l'agent) : seules les colonnes
            # numériques sont sommées — un total faux est impossible par construction.
            total = ["TOTAL"]
            for idx, col in enumerate(colonnes[1:], start=1):
                valeurs = [l[idx] for l in lignes if len(l) > idx and isinstance(l[idx], (int, float))]
                total.append(round(sum(valeurs), 2) if valeurs and str(col.get("format")) in ("euro", "nombre") else "")
            ws.append(total)
            for cellule in ws[ws.max_row]:
                cellule.font = Font(bold=True)
            for idx, col in enumerate(colonnes, start=1):
                fmt = formats.get(str(col.get("format", "texte")))
                if fmt:
                    ws.cell(row=ws.max_row, column=idx).number_format = fmt
    wb.save(args.sortie)
    print(f"Tableur créé : {args.sortie}")


# ── markdown minimal → blocs (titres, puces, paragraphes) ────────────────────


def _blocs_markdown(md: str) -> list[tuple[str, str]]:
    blocs: list[tuple[str, str]] = []
    for ligne in md.splitlines():
        brut = ligne.strip()
        if not brut:
            continue
        if brut.startswith("### "):
            blocs.append(("h3", brut[4:]))
        elif brut.startswith("## "):
            blocs.append(("h2", brut[3:]))
        elif brut.startswith("# "):
            blocs.append(("h1", brut[2:]))
        elif brut.startswith(("- ", "* ")):
            blocs.append(("puce", brut[2:]))
        else:
            blocs.append(("p", brut))
    return blocs


def cmd_docx(args: argparse.Namespace) -> None:
    import docx
    from docx.shared import Pt, RGBColor

    document = docx.Document()
    accent = RGBColor(0x1F, 0x38, 0x64)
    titre = document.add_heading(args.titre, level=0)
    for run in titre.runs:
        run.font.color.rgb = accent
    niveaux = {"h1": 1, "h2": 2, "h3": 3}
    for genre, texte in _blocs_markdown(_lire(args.md)):
        if genre in niveaux:
            entete = document.add_heading(texte, level=niveaux[genre])
            for run in entete.runs:
                run.font.color.rgb = accent
        elif genre == "puce":
            document.add_paragraph(texte, style="List Bullet")
        else:
            paragraphe = document.add_paragraph(texte)
            for run in paragraphe.runs:
                run.font.size = Pt(11)
    document.save(args.sortie)
    print(f"Document Word créé : {args.sortie}")


def cmd_pdf(args: argparse.Namespace) -> None:
    import markdown as md_lib
    from xhtml2pdf import pisa

    corps = md_lib.markdown(_lire(args.md), extensions=["tables"])
    html = f"""<html><head><style>
      @page {{ size: A4; margin: 2cm; }}
      body {{ font-family: Helvetica, sans-serif; font-size: 11pt; color: #1a1a1a; }}
      h1 {{ color: #{BLEU_NUIT}; border-bottom: 2px solid #{BLEU_NUIT}; padding-bottom: 4px; }}
      h2, h3 {{ color: #{BLEU_NUIT}; }}
      table {{ border-collapse: collapse; width: 100%; margin: 8px 0; }}
      th {{ background-color: #{BLEU_NUIT}; color: white; padding: 6px; text-align: left; }}
      td {{ border: 1px solid #cccccc; padding: 6px; }}
    </style></head><body><h1>{args.titre}</h1>{corps}</body></html>"""
    with open(args.sortie, "wb") as sortie:
        resultat = pisa.CreatePDF(html, dest=sortie, encoding="utf-8")
    if resultat.err:
        _fail("La génération PDF a échoué (contenu markdown à simplifier).")
    print(f"PDF créé : {args.sortie}")


# ── pptx : présentation sobre (les slides de haute qualité = MCP presenton) ──
# spec.json : {"titre": "...", "sous_titre": "...", "slides": [{"titre": "...", "points": ["..."]}]}


def cmd_pptx(args: argparse.Namespace) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    spec = _lire_json(args.spec)
    accent = RGBColor(0x1F, 0x38, 0x64)
    presentation = Presentation()
    garde = presentation.slides.add_slide(presentation.slide_layouts[0])
    garde.shapes.title.text = spec.get("titre") or "Présentation"
    garde.shapes.title.text_frame.paragraphs[0].font.color.rgb = accent
    if spec.get("sous_titre"):
        garde.placeholders[1].text = spec["sous_titre"]
    for slide_spec in spec.get("slides") or []:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = slide_spec.get("titre", "")
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = accent
        corps = slide.placeholders[1].text_frame
        corps.clear()
        for rang, point in enumerate(slide_spec.get("points") or []):
            paragraphe = corps.paragraphs[0] if rang == 0 else corps.add_paragraph()
            paragraphe.text = str(point)
            paragraphe.font.size = Pt(18)
    presentation.save(args.sortie)
    print(f"Présentation créée : {args.sortie}")


# ── publier : pont Fichiers (API officielle, clé interne) ────────────────────


def cmd_publier(args: argparse.Namespace) -> None:
    api_key = os.environ.get("LUNARIA_INTERNAL_API_KEY", "").strip()
    base_url = os.environ.get("LUNARIA_APP_URL", DEFAULT_APP_URL).rstrip("/")
    if not api_key:
        _fail("Clé interne absente (LUNARIA_INTERNAL_API_KEY) : pont Fichiers non configuré.")
    chemin = args.fichier
    if not os.path.isfile(chemin):
        _fail(f"Fichier introuvable : {chemin}")
    nom = args.nom or os.path.basename(chemin)
    ctype = mimetypes.guess_type(nom)[0] or "application/octet-stream"
    frontiere = uuid.uuid4().hex
    with open(chemin, "rb") as handle:
        contenu = handle.read()
    corps = (
        f"--{frontiere}\r\nContent-Disposition: form-data; name=\"file\"; filename=\"{nom}\"\r\n"
        f"Content-Type: {ctype}\r\n\r\n"
    ).encode("utf-8") + contenu + f"\r\n--{frontiere}--\r\n".encode("utf-8")
    # process=false : un document généré n'a pas à être indexé pour la recherche.
    req = urllib.request.Request(f"{base_url}/api/v1/files/?process=false", data=corps, method="POST")
    req.add_header("Authorization", f"Bearer {api_key}")
    req.add_header("Content-Type", f"multipart/form-data; boundary={frontiere}")
    try:
        with urllib.request.urlopen(req, timeout=TIMEOUT_SECONDS) as resp:
            fichier = json.loads(resp.read().decode("utf-8"))
    except urllib.error.HTTPError as exc:
        _fail(f"L'application a refusé le téléversement ({exc.code}). {exc.read().decode('utf-8', 'replace')[:200]}")
    except urllib.error.URLError as exc:
        _fail(f"Application injoignable ({exc.reason}).")
    identifiant = fichier.get("id")
    if not identifiant:
        _fail("Téléversement sans identifiant retourné.")
    print(f"Fichier publié : {nom}")
    print(f"Lien de téléchargement à donner au patron : [{nom}](/api/v1/files/{identifiant}/content)")


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(prog="doc_cli", description="Atelier documents LunarIA (Max).")
    sub = parser.add_subparsers(dest="command", required=True)

    p_xlsx = sub.add_parser("xlsx", help="Tableur Excel sobre, totaux calculés.")
    p_xlsx.add_argument("--titre", required=True)
    p_xlsx.add_argument("--data-file", required=True, help="Spec JSON (feuilles/colonnes/lignes).")
    p_xlsx.add_argument("--sortie", required=True)
    p_xlsx.set_defaults(func=cmd_xlsx)

    p_docx = sub.add_parser("docx", help="Document Word depuis markdown simple.")
    p_docx.add_argument("--titre", required=True)
    p_docx.add_argument("--md", required=True, help="Fichier markdown (titres, puces, paragraphes).")
    p_docx.add_argument("--sortie", required=True)
    p_docx.set_defaults(func=cmd_docx)

    p_pdf = sub.add_parser("pdf", help="Rapport PDF depuis markdown (tableaux acceptés).")
    p_pdf.add_argument("--titre", required=True)
    p_pdf.add_argument("--md", required=True)
    p_pdf.add_argument("--sortie", required=True)
    p_pdf.set_defaults(func=cmd_pdf)

    p_pptx = sub.add_parser("pptx", help="Présentation sobre (haute qualité = MCP presenton).")
    p_pptx.add_argument("--spec", required=True, help="Spec JSON (titre, sous_titre, slides).")
    p_pptx.add_argument("--sortie", required=True)
    p_pptx.set_defaults(func=cmd_pptx)

    p_pub = sub.add_parser("publier", help="Téléverse le document et rend le lien de téléchargement.")
    p_pub.add_argument("--fichier", required=True)
    p_pub.add_argument("--nom", help="Nom affiché (défaut : nom du fichier).")
    p_pub.set_defaults(func=cmd_publier)

    return parser


def main() -> None:
    args = build_parser().parse_args()
    args.func(args)


if __name__ == "__main__":
    main()
